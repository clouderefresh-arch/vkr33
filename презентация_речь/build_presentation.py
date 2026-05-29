# -*- coding: utf-8 -*-
"""
Сборка презентации к ВКР Кострова на основе шаблона «шаблон-СПО.pptx».

Скрипт заполняет 10-слайдовую структуру шаблона данными из дипломной работы
«Применение интернет-технологий в коммерческой деятельности предприятий».

Принцип: структура и оформление шаблона сохраняются, текст-заполнители
заменяются реальными данными ВКР. Брендинг примера («Гемотест») удаляется.
"""
import copy
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

SRC = "исходники/шаблон-СПО.pptx"
OUT = "Презентация_Костров_ВКР.pptx"

NAVY = RGBColor(0x00, 0x20, 0x60)
DARK = RGBColor(0x22, 0x22, 0x22)
GREEN = RGBColor(0x1F, 0x6E, 0x43)

prs = Presentation(SRC)


def shape_by_name(slide, name):
    for sh in slide.shapes:
        if sh.name == name:
            return sh
    return None


def set_runs_text(shape, text):
    """Заменить текст в первом run первого абзаца, очистить остальные."""
    tf = shape.text_frame
    p0 = tf.paragraphs[0]
    if p0.runs:
        p0.runs[0].text = text
        for r in p0.runs[1:]:
            r._r.getparent().remove(r._r)
    else:
        p0.add_run().text = text
    # удалить лишние абзацы
    for p in tf.paragraphs[1:]:
        p._p.getparent().remove(p._p)


def fill_textframe(tf, items, default_size=14, default_color=DARK,
                   font_name="Arial", space_after=4):
    """Полностью пересобрать текстовый фрейм из списка строк.

    items: список словарей {text, size, bold, color, level, align}.
    """
    # очистить все абзацы кроме первого
    for p in list(tf.paragraphs[1:]):
        p._p.getparent().remove(p._p)
    first = tf.paragraphs[0]
    for r in list(first.runs):
        r._r.getparent().remove(r._r)

    for i, it in enumerate(items):
        p = first if i == 0 else tf.add_paragraph()
        p.level = it.get("level", 0)
        p.alignment = it.get("align", PP_ALIGN.LEFT)
        p.space_after = Pt(it.get("space_after", space_after))
        p.space_before = Pt(it.get("space_before", 0))
        run = p.add_run()
        run.text = it["text"]
        f = run.font
        f.name = it.get("font", font_name)
        f.size = Pt(it.get("size", default_size))
        f.bold = it.get("bold", False)
        f.color.rgb = it.get("color", default_color)


def add_title(slide, text, top=0.5, size=28, color=NAVY):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(9.0), Inches(0.9))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = "Century Gothic"
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = color
    return tb


def add_body(slide, items, left=0.6, top=1.6, width=8.9, height=5.4,
             default_size=15):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    fill_textframe(tf, items, default_size=default_size)
    return tb


def style_table(table, header_size=11, body_size=11):
    for ci, cell in enumerate(table.rows[0].cells):
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.size = Pt(header_size)
                r.font.bold = True
                r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                r.font.name = "Arial"
    for ri, row in enumerate(table.rows):
        if ri == 0:
            continue
        for cell in row.cells:
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(body_size)
                    r.font.name = "Arial"
                    r.font.color.rgb = DARK


def add_table(slide, data, left, top, width, height, col_widths=None,
              header_size=11, body_size=11):
    rows, cols = len(data), len(data[0])
    gfx = slide.shapes.add_table(rows, cols, Inches(left), Inches(top),
                                 Inches(width), Inches(height))
    table = gfx.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
    for ri, row in enumerate(data):
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = str(val)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_top = Pt(2)
            cell.margin_bottom = Pt(2)
            cell.margin_left = Pt(4)
            cell.margin_right = Pt(4)
    style_table(table, header_size, body_size)
    return table


# ---------------------------------------------------------------------------
# СЛАЙД 1 — Титульный
# ---------------------------------------------------------------------------
s = prs.slides[0]
topic = shape_by_name(s, "TextBox 3")
set_runs_text(topic, "«Применение интернет-технологий в коммерческой деятельности предприятий»")

rect = shape_by_name(s, "Прямоугольник 5")
# абзац 0 — руководитель (оставляем), абзац 1 — студент
tf = rect.text_frame
# найдём абзац со «Студент»
for p in tf.paragraphs:
    joined = "".join(r.text for r in p.runs)
    if "Студент" in joined:
        if p.runs:
            p.runs[0].text = "Студент: Костров"
            for r in p.runs[1:]:
                r._r.getparent().remove(r._r)

# ---------------------------------------------------------------------------
# СЛАЙД 2 — Актуальность, объект, предмет, цель, задачи
# ---------------------------------------------------------------------------
s = prs.slides[1]
title = shape_by_name(s, "TextBox 2")
set_runs_text(title, "Актуальность, объект, предмет, цель и задачи")

body = shape_by_name(s, "Прямоугольник 4")
body.top = Inches(1.45)
body.left = Inches(0.45)
body.width = Inches(9.2)
body.height = Inches(5.7)
items = [
    {"text": "Актуальность: рынок интернет-торговли РФ в 2024 г. достиг ≈ 9,1 трлн руб. "
             "(+35 % к 2023 г.), доля онлайн в рознице — 15–16 %. Рост цифровых каналов "
             "усиливает угрозы экономической безопасности коммерческих предприятий.",
     "size": 13, "space_after": 6},
    {"text": "Объект — коммерческие организации ЮЗАО г. Москвы, ведущие торгово-сбытовую "
             "деятельность с использованием интернет-технологий (модельная группа П-1 — П-5).",
     "size": 13, "space_after": 4},
    {"text": "Предмет — перспективы применения интернет-технологий и отношения предприятий "
             "с подразделением экономической безопасности по устойчивости цифровых каналов продаж.",
     "size": 13, "space_after": 4},
    {"text": "Цель — разработать комплекс концептуальных предложений по повышению "
             "эффективности коммерческой деятельности за счёт усиления экономической "
             "безопасности при поддержке ОЭБиПК УВД по ЮЗАО.",
     "size": 13, "bold": True, "color": NAVY, "space_after": 6},
    {"text": "Задачи:", "size": 13, "bold": True, "space_after": 2},
    {"text": "• раскрыть понятийно-сущностные основы интернет-технологий в коммерции;", "size": 12, "level": 1, "space_after": 1},
    {"text": "• систематизировать зарубежный и отечественный опыт цифровых экосистем;", "size": 12, "level": 1, "space_after": 1},
    {"text": "• определить роль экономической безопасности как фактора эффективности;", "size": 12, "level": 1, "space_after": 1},
    {"text": "• дать организационно-правовую характеристику ОЭБиПК УВД по ЮЗАО;", "size": 12, "level": 1, "space_after": 1},
    {"text": "• проанализировать модельную группу предприятий и их цифровые каналы;", "size": 12, "level": 1, "space_after": 1},
    {"text": "• выявить уязвимости и угрозы экономической безопасности;", "size": 12, "level": 1, "space_after": 1},
    {"text": "• разработать программу мер (профилактика, аналитика, мониторинг, обмен данными);", "size": 12, "level": 1, "space_after": 1},
    {"text": "• оценить экономическую эффективность и определить перспективы развития.", "size": 12, "level": 1, "space_after": 1},
]
fill_textframe(body.text_frame, items)
body.text_frame.word_wrap = True

# ---------------------------------------------------------------------------
# СЛАЙД 3 — База исследования (ОЭБиПК + модельная группа)
# ---------------------------------------------------------------------------
s = prs.slides[2]
rect = shape_by_name(s, "Прямоугольник 2")
rect.top = Inches(0.5)
rect.left = Inches(0.5)
rect.width = Inches(9.0)
rect.height = Inches(0.9)
fill_textframe(rect.text_frame, [
    {"text": "База исследования: ОЭБиПК УВД по ЮЗАО и модельная группа предприятий",
     "size": 24, "bold": True, "color": NAVY, "font": "Century Gothic"}
])
rect.text_frame.word_wrap = True
add_body(s, [
    {"text": "ОЭБиПК УВД по ЮЗАО ГУ МВД России по г. Москве — подразделение экономической "
             "безопасности и противодействия коррупции. Не является коммерческой организацией, "
             "но играет ключевую инфраструктурную роль для устойчивости бизнеса округа; его "
             "функция смещается от реагирования к профилактике и цифровому мониторингу.",
     "size": 14, "space_after": 8},
    {"text": "Модельная группа — 5 предприятий-аналогов ЮЗАО (обезличены, ст. 161 УПК РФ, ФЗ № 152-ФЗ):",
     "size": 14, "bold": True, "space_after": 3},
    {"text": "• П-1 — инжиниринг и торговля оборудованием (выручка 146,4 млн руб., онлайн 31,5 %);", "size": 13, "level": 1, "space_after": 1},
    {"text": "• П-2 — частный медицинский центр (37,2 млн руб., онлайн 22,8 %);", "size": 13, "level": 1, "space_after": 1},
    {"text": "• П-3 — оптово-розничная торговля стройматериалами (284,1 млн руб., онлайн 18,4 %);", "size": 13, "level": 1, "space_after": 1},
    {"text": "• П-4 — фермерский маркетплейс (62,8 млн руб., онлайн 44,7 %);", "size": 13, "level": 1, "space_after": 1},
    {"text": "• П-5 — стоматологическая сеть (98,7 млн руб., онлайн 19,6 %).", "size": 13, "level": 1, "space_after": 1},
], top=1.55, height=5.5)

# ---------------------------------------------------------------------------
# СЛАЙД 4 — Цветной слайд «в цифрах» (удаляем брендинг «Гемотест»)
# ---------------------------------------------------------------------------
s = prs.slides[3]
# удалить картинки-логотипы и фоновое фото
for sh in list(s.shapes):
    if sh.shape_type == 13:  # PICTURE
        sh._element.getparent().remove(sh._element)
# заголовок
head = shape_by_name(s, "TextBox 4")
fill_textframe(head.text_frame, [
    {"text": "Цифровая среда коммерческой деятельности и модельная группа — в цифрах",
     "size": 24, "bold": True, "color": NAVY, "font": "Arial Black"}
])
head.text_frame.word_wrap = True
head.top = Inches(1.4)
# нижняя подпись-источник
note = shape_by_name(s, "Прямоугольник 2")
if note is not None:
    fill_textframe(note.text_frame, [
        {"text": "По данным АКИТ, «Дата Инсайт», Росстата и расчётам автора (2024 г.)",
         "size": 11, "color": RGBColor(0x80, 0x80, 0x80), "font": "Arial"}
    ])

card_content = {
    "Скругленный прямоугольник 16": ("≈ 9,1 трлн ₽", "объём рынка интернет-торговли России, 2024"),
    "Скругленный прямоугольник 17": ("+35 %", "годовой рост рынка к 2023 году"),
    "Скругленный прямоугольник 18": ("15–16 %", "доля онлайн-канала в обороте розницы"),
    "Скругленный прямоугольник 19": ("5", "предприятий модельной группы (П-1 — П-5)"),
    "Скругленный прямоугольник 21": ("0,43–0,86", "диапазон индекса цифрового доверия (ИЦД)"),
    "Скругленный прямоугольник 20": ("KEO ≈ 9,5", "отдача на рубль затрат (базовый сценарий)"),
}
for name, (big, label) in card_content.items():
    sh = shape_by_name(s, name)
    if sh is None:
        continue
    sh.fill.solid()
    sh.fill.fore_color.rgb = RGBColor(0x2F, 0x55, 0x97)
    sh.line.fill.background()
    fill_textframe(sh.text_frame, [
        {"text": big, "size": 18, "bold": True, "color": RGBColor(0xFF, 0xFF, 0xFF),
         "align": PP_ALIGN.CENTER, "space_after": 2},
        {"text": label, "size": 9.5, "bold": True, "color": RGBColor(0xFF, 0xFF, 0xFF),
         "align": PP_ALIGN.CENTER, "space_after": 0},
    ])
    sh.text_frame.word_wrap = True

# ---------------------------------------------------------------------------
# СЛАЙД 5 — Анализ цифровых каналов модельной группы (глава 2)
# ---------------------------------------------------------------------------
s = prs.slides[4]
rect = shape_by_name(s, "Прямоугольник 2")
rect.top = Inches(0.5); rect.left = Inches(0.5)
rect.width = Inches(9.0); rect.height = Inches(0.8)
fill_textframe(rect.text_frame, [
    {"text": "Анализ цифровых каналов модельной группы (глава 2)",
     "size": 24, "bold": True, "color": NAVY, "font": "Century Gothic"}
])
rect.text_frame.word_wrap = True
add_table(s, [
    ["Параметр", "П-1", "П-2", "П-3", "П-4", "П-5"],
    ["Профиль", "Инжиниринг/\nоборуд.", "Медцентр", "Строймат.", "Фермер.\nмаркетпл.", "Стомат.\nсеть"],
    ["Выручка 2024, млн ₽", "146,4", "37,2", "284,1", "62,8", "98,7"],
    ["Доля онлайн, %", "31,5", "22,8", "18,4", "44,7", "19,6"],
    ["Заказов онлайн/мес", "1 240", "2 870", "1 880", "9 350", "4 110"],
    ["Средний чек, ₽", "31 600", "3 950", "12 470", "1 750", "5 320"],
    ["Коэф. цифр. интенс. Kди", "0,132", "0,146", "0,057", "0,367", "0,108"],
], left=0.5, top=1.5, width=9.0, height=4.4,
   col_widths=[2.4, 1.32, 1.32, 1.32, 1.32, 1.32], header_size=12, body_size=11)
add_body(s, [
    {"text": "Все предприятия используют национальную инфраструктуру доверия (СБП, эквайринг, "
             "маркировка «Честный знак»), но различаются по цифровой интенсивности и устойчивости.",
     "size": 12, "color": DARK}
], left=0.5, top=6.1, width=9.0, height=1.0, default_size=12)

# ---------------------------------------------------------------------------
# СЛАЙД 6 — Защищённость каналов и индекс цифрового доверия (ИЦД)
# ---------------------------------------------------------------------------
s = prs.slides[5]
rect = shape_by_name(s, "Прямоугольник 4")
rect.top = Inches(0.5); rect.left = Inches(0.5)
rect.width = Inches(9.0); rect.height = Inches(0.8)
fill_textframe(rect.text_frame, [
    {"text": "Защищённость каналов и индекс цифрового доверия (ИЦД)",
     "size": 24, "bold": True, "color": NAVY, "font": "Century Gothic"}
])
rect.text_frame.word_wrap = True
add_table(s, [
    ["Предприятие", "П-1", "П-2", "П-3", "П-4", "П-5"],
    ["Средний уровень мер защиты, %", "78,6", "42,9", "71,4", "85,7", "50,0"],
    ["Индекс цифрового доверия (ИЦД)", "0,71", "0,43", "0,64", "0,86", "0,50"],
    ["Зона устойчивости", "Приемл.", "Ниже\nсред.", "Приемл.", "Высокий", "Ниже\nсред."],
], left=0.5, top=1.5, width=9.0, height=2.6,
   col_widths=[3.0, 1.2, 1.2, 1.2, 1.2, 1.2], header_size=12, body_size=11)
add_body(s, [
    {"text": "Авторский индекс цифрового доверия (ИЦД) рассчитан по 7 индикаторам защиты "
             "(шкала 0–1).", "size": 13, "space_after": 4},
    {"text": "• ИЦД положительно коррелирует с долей выручки от онлайн-канала;", "size": 12, "level": 1, "space_after": 1},
    {"text": "• самая распространённая проблема — отсутствие обучения персонала кибергигиене "
             "(в среднем 20 %);", "size": 12, "level": 1, "space_after": 1},
    {"text": "• наиболее уязвимы П-2 и П-5 (медицинский и стоматологический профиль).", "size": 12, "level": 1, "space_after": 1},
], left=0.5, top=4.4, width=9.0, height=2.7, default_size=13)

# ---------------------------------------------------------------------------
# СЛАЙД 7 — Ключевые уязвимости и угрозы экономической безопасности
# ---------------------------------------------------------------------------
s = prs.slides[6]
rect = shape_by_name(s, "Прямоугольник 2")
rect.top = Inches(0.5); rect.left = Inches(0.5)
rect.width = Inches(9.0); rect.height = Inches(0.8)
fill_textframe(rect.text_frame, [
    {"text": "Ключевые уязвимости и угрозы экономической безопасности",
     "size": 24, "bold": True, "color": NAVY, "font": "Century Gothic"}
])
rect.text_frame.word_wrap = True
add_body(s, [
    {"text": "Систематизированы характерные для предприятий ЮЗАО уязвимости:", "size": 15, "bold": True, "space_after": 5},
    {"text": "• отсутствие или фрагментарность программ обучения персонала кибергигиене;", "size": 14, "level": 1, "space_after": 3},
    {"text": "• неполная защита платёжных операций антифрод-модулями;", "size": 14, "level": 1, "space_after": 3},
    {"text": "• пропуски в работе с цифровой маркировкой «Честный знак»;", "size": 14, "level": 1, "space_after": 3},
    {"text": "• высокая восприимчивость медицинского и стоматологического профиля к фишингу;", "size": 14, "level": 1, "space_after": 3},
    {"text": "• рост доли оспоренных платежей (chargeback) и возвратов;", "size": 14, "level": 1, "space_after": 3},
    {"text": "• появление фишинговых копий сайтов и приложений предприятий;", "size": 14, "level": 1, "space_after": 3},
    {"text": "• медленная реакция службы поддержки на обращения и инциденты.", "size": 14, "level": 1, "space_after": 3},
], top=1.55, height=5.4, default_size=14)

# ---------------------------------------------------------------------------
# СЛАЙД 8 — Концептуальные предложения: программа из 4 звеньев
# ---------------------------------------------------------------------------
s = prs.slides[7]
rect = shape_by_name(s, "Прямоугольник 2")
rect.top = Inches(0.5); rect.left = Inches(0.5)
rect.width = Inches(9.0); rect.height = Inches(0.8)
fill_textframe(rect.text_frame, [
    {"text": "Концептуальные предложения: программа из 4 звеньев",
     "size": 24, "bold": True, "color": NAVY, "font": "Century Gothic"}
])
rect.text_frame.word_wrap = True
add_table(s, [
    ["Звено", "Инструменты", "Эффект"],
    ["1. Профилактика", "Семинары, курс «Кибергигиена для МСП», чек-листы",
     "Снижение числа инцидентов до их наступления"],
    ["2. Аналитика", "Обезличенные ежеквартальные обзоры типовых схем",
     "Обновление внутренних регламентов предприятий"],
    ["3. Цифровой мониторинг", "АИС/ЕИС МВД, СПАРК, Контур.Фокус, ФинЦЕРТ; 5 индикаторов",
     "Раннее выявление подозрительных операций"],
    ["4. Межведомственный обмен", "Рабочая группа бизнес—ЭБиПК—ФНС—Банк России; ГосТех",
     "Целостная картина рисков, синхронизация мер"],
], left=0.5, top=1.5, width=9.0, height=3.6,
   col_widths=[2.3, 3.6, 3.1], header_size=12, body_size=11)
add_body(s, [
    {"text": "7 конкретных мероприятий; программа укладывается в действующие полномочия ОЭБиПК "
             "(ФЗ № 3-ФЗ «О полиции», ФЗ № 273-ФЗ) без расширения функций. Реализация — в 3 фазы.",
     "size": 12, "color": DARK}
], left=0.5, top=5.4, width=9.0, height=1.4, default_size=12)

# ---------------------------------------------------------------------------
# СЛАЙД 9 — Экономическая эффективность (вывод)
# ---------------------------------------------------------------------------
s = prs.slides[8]
rect = shape_by_name(s, "Прямоугольник 2")
rect.top = Inches(0.5); rect.left = Inches(0.5)
rect.width = Inches(9.0); rect.height = Inches(0.8)
fill_textframe(rect.text_frame, [
    {"text": "Экономическая эффективность предложений (вывод)",
     "size": 24, "bold": True, "color": NAVY, "font": "Century Gothic"}
])
rect.text_frame.word_wrap = True
add_table(s, [
    ["Показатель (среднее по группе)", "2024 (факт)", "2026 (прогноз)", "Изменение"],
    ["Выручка от онлайн-канала, млн ₽", "48,3", "63,7", "+31,9 %"],
    ["Конверсия сайта в покупку, %", "1,9", "3,1", "+1,2 п.п."],
    ["Потери от фрода и фишинга, млн ₽/год", "1,42", "0,28", "−80,3 %"],
    ["Индекс лояльности клиентов", "28", "51", "+23 п."],
    ["Чистая прибыль, млн ₽/год", "4,7", "8,9", "+89,4 %"],
], left=0.5, top=1.45, width=9.0, height=3.0,
   col_widths=[4.0, 1.7, 1.7, 1.6], header_size=12, body_size=11)
add_body(s, [
    {"text": "Интегральная оценка эффективности:", "size": 14, "bold": True, "space_after": 4},
    {"text": "• Коэффициент экономической отдачи KEO = (4 200 − 400) / 400 = 9,5 "
             "(≈ 10 ₽ эффекта на 1 ₽ затрат);", "size": 13, "level": 1, "space_after": 2},
    {"text": "• NPV ≈ 9 250 тыс. ₽ при инвестициях ≈ 150 тыс. ₽/предприятие; срок окупаемости — 5,3 мес.;", "size": 13, "level": 1, "space_after": 2},
    {"text": "• дополнительный налоговый эффект по группе — 4,9–5,4 млн ₽/год.", "size": 13, "level": 1, "space_after": 2},
    {"text": "Цель работы достигнута, гипотеза подтверждена.", "size": 13, "bold": True, "color": GREEN, "space_before": 4},
], left=0.5, top=4.55, width=9.0, height=2.6, default_size=13)

# СЛАЙД 10 — «Спасибо за внимание» оставляем без изменений.

prs.save(OUT)
print("Saved:", OUT, "slides:", len(prs.slides))
